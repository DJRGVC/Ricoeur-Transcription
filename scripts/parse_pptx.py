from pptx import Presentation
import os
import shutil

# Define the directory to store the images and text
save_dir = "extracted_data"
images_dir = os.path.join(save_dir, "images")
text_file = os.path.join(save_dir, "slides_text.txt")

# Create directories for storing images and text
if os.path.exists(save_dir):
    shutil.rmtree(save_dir)
os.makedirs(images_dir)

# Load the presentation
prs = Presentation("your_presentation.pptx")  # Change this to the path of your downloaded presentation

# Open text file to write slide text
with open(text_file, 'w') as text_output:

    # Iterate over slides
    for idx, slide in enumerate(prs.slides):
        slide_text = ""
        image_idx = 1
        
        # Extract text
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text.strip() + "\n"
        
        # Write slide text to file
        text_output.write(f"Slide {idx + 1}:\n{slide_text}\n{'-'*20}\n")
        
        # Extract images
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_file = os.path.join(images_dir, f"slide_{idx + 1}_image_{image_idx}.png")
                with open(image_file, "wb") as img_out:
                    img_out.write(image.blob)
                image_idx += 1
                
print(f"Text extracted to {text_file}")
print(f"Images extracted to {images_dir}")

